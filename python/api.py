#!/usr/bin/env python3
# -*- coding: utf-8 -*-

import os
import sys
import json
import time
import subprocess
import shutil
import traceback
import urllib.parse
from typing import List, Dict, Any, Optional

# 配置详细日志
import logging

# 创建自定义的日志格式和处理器
class StdoutLogHandler(logging.StreamHandler):
    def __init__(self):
        logging.StreamHandler.__init__(self, sys.stdout)

# 创建控制台处理器，正常日志输出到标准输出
stdout_handler = StdoutLogHandler()
stdout_handler.setLevel(logging.INFO)
stdout_handler.setFormatter(logging.Formatter('%(asctime)s - %(name)s - %(levelname)s - %(message)s'))

# 配置根日志记录器
root_logger = logging.getLogger()
root_logger.setLevel(logging.INFO)
# 移除所有已存在的处理器
for handler in root_logger.handlers[:]:
    root_logger.removeHandler(handler)
root_logger.addHandler(stdout_handler)

# 获取应用程序日志记录器
logger = logging.getLogger(__name__)

# 调整第三方库的日志级别
logging.getLogger("faiss").setLevel(logging.ERROR)
logging.getLogger("langchain").setLevel(logging.WARNING)
logging.getLogger("transformers").setLevel(logging.ERROR)
logging.getLogger("huggingface_hub").setLevel(logging.ERROR)
logging.getLogger("dashscope").setLevel(logging.ERROR)
logging.getLogger("httpx").setLevel(logging.ERROR)
logging.getLogger("urllib3").setLevel(logging.ERROR)

# 检查和加载环境变量
def load_env_variables():
    """不再加载环境变量，直接使用硬编码配置"""
    print("配置已硬编码在程序中，不再从环境变量加载")
    return True

# 在导入其他模块前加载环境变量
load_env_variables()

# 导入watchdog相关模块
from watchdog.observers import Observer
from watchdog.events import FileSystemEventHandler, FileCreatedEvent, FileDeletedEvent, FileMovedEvent, FileModifiedEvent


from fastapi import FastAPI, BackgroundTasks, Request
from fastapi.middleware.cors import CORSMiddleware
from fastapi.responses import JSONResponse
from pydantic import BaseModel
import uvicorn

# 导入所需的langchain模块
from langchain_community.document_loaders import (
    TextLoader, 
    PyPDFLoader
)
from langchain.document_loaders.base import BaseLoader
from langchain.schema import Document
from langchain.text_splitter import CharacterTextSplitter
from langchain_community.vectorstores import FAISS
from langchain_community.embeddings import DashScopeEmbeddings
from dotenv import load_dotenv

# 全局变量：文件监控相关
active_observers = {}  # 存储活跃的文件监控器
file_handlers = {}  # 存储文件处理器

# 索引状态全局变量
index_status = {
    "in_progress": False,
    "progress": 0,
    "status": "",
    "completed": False,
    "error": None,
    "file_stats": {
        "success_count": 0,
        "failure_count": 0,
        "skipped_count": 0,
        "total_count": 0
    },
    "success_files": [],    # 成功索引的文件列表
    "failed_files": [],     # 索引失败的文件列表
    "skipped_files": []     # 跳过的文件列表
}

# 全局配置参数
MAX_TEXT_LENGTH = 20000  # 每个文件的最大字符数限制
MAX_CHUNK_COUNT = 200    # 每个文件最大分块数量
MAX_FILE_SIZE_MB = 100   # 默认最大文件大小限制(MB)，从10MB改为100MB
MAX_TEXT_BLOCK_SIZE = 2048  # v2模型限制为2048 Token/行
MAX_BATCH_ROWS = 25      # 通义千问一次调用支持的最大行数

# 直接设置通义千问API密钥
# 请替换成你自己的通义千问API密钥
DASHSCOPE_API_KEY = "your_api_key_here"

# 嵌入模型配置
EMBEDDING_MODELS = {
    "text-embedding-v2": {
        "name": "文本嵌入模型V2",
        "max_tokens": 2048,
        "supported_languages": "中文、英文、多语言支持"
    },
    "text-embedding-v3": {
        "name": "文本嵌入模型V3",
        "max_tokens": 8192,
        "supported_languages": "中文、英文、多语言支持(50+语种)"
    }
}

# 默认使用v2模型
EMBEDDING_MODEL_NAME = "text-embedding-v2"

# 根据当前选择的模型自动更新token限制
def update_token_limit():
    global MAX_TEXT_BLOCK_SIZE
    if EMBEDDING_MODEL_NAME in EMBEDDING_MODELS:
        MAX_TEXT_BLOCK_SIZE = EMBEDDING_MODELS[EMBEDDING_MODEL_NAME]["max_tokens"]
        logger.info(f"已更新文本块大小限制为 {MAX_TEXT_BLOCK_SIZE} (基于模型 {EMBEDDING_MODEL_NAME})")

# 确保token限制是正确的
update_token_limit()

# 导入其他模块前调用load_env_variables函数（现在只是一个保留的函数调用）
load_env_variables()

# 文件监控类
class FileIndexHandler(FileSystemEventHandler):
    def __init__(self, folder_path):
        self.folder_path = folder_path
        self.db_path = get_db_path(folder_path)
        self.pending_events = []
        self.processing_lock = threading.Lock()
        self.last_update_time = 0
        self.update_timer = None
        
        # 启动处理线程
        self.update_thread = threading.Thread(target=self.process_events_thread, daemon=True)
        self.update_thread.start()
        logger.info(f"初始化文件监控: {folder_path}")
    
    def on_created(self, event):
        """当文件或目录被创建时"""
        if not event.is_directory and is_supported_file(event.src_path):
            logger.info(f"文件创建: {event.src_path}")
            self.schedule_update(event)
    
    def on_deleted(self, event):
        """当文件或目录被删除时"""
        if not event.is_directory:
            # 这里不用is_supported_file检查，因为文件已经被删除
            # 但我们可以检查文件名是否应该被忽略
            file_name = os.path.basename(event.src_path)
            if (file_name.startswith('.') or file_name in ['.DS_Store', 'Thumbs.db', 'desktop.ini'] or 
                file_name.endswith('~') or file_name.startswith('~$')):
                return  # 忽略系统隐藏文件的删除事件
                
            logger.info(f"文件删除: {event.src_path}")
            self.schedule_update(event)
    
    def on_moved(self, event):
        """当文件或目录被移动或重命名时"""
        if not event.is_directory:
            # 源路径可能已不存在，但检查目标路径
            file_name_src = os.path.basename(event.src_path)
            if (file_name_src.startswith('.') or file_name_src in ['.DS_Store', 'Thumbs.db', 'desktop.ini'] or 
                file_name_src.endswith('~') or file_name_src.startswith('~$')):
                return  # 忽略系统隐藏文件
                
            # 将其视为删除旧文件并创建新文件
            logger.info(f"文件移动: {event.src_path} -> {event.dest_path}")
            self.schedule_update(event)
    
    def on_modified(self, event):
        """当文件被修改时"""
        if not event.is_directory and is_supported_file(event.src_path):
            logger.info(f"文件修改: {event.src_path}")
            self.schedule_update(event)
    
    def schedule_update(self, event):
        """安排更新，延迟执行以避免频繁更新"""
        with self.processing_lock:
            self.pending_events.append(event)
            
            # 取消之前计划的更新
            if self.update_timer:
                self.update_timer.cancel()
            
            # 安排新的更新，延迟5秒
            self.update_timer = threading.Timer(5.0, self.trigger_update)
            self.update_timer.daemon = True
            self.update_timer.start()
    
    def trigger_update(self):
        """触发更新，通知处理线程处理事件"""
        logger.info("触发索引更新")
        self.last_update_time = time.time()
    
    def process_events_thread(self):
        """事件处理线程，定期检查是否有待处理的事件"""
        while True:
            if time.time() - self.last_update_time >= 5 and self.pending_events:
                with self.processing_lock:
                    if self.pending_events:
                        # 收集所有待处理事件
                        events = self.pending_events.copy()
                        self.pending_events.clear()
                        # 处理事件
                        self.process_events(events)
            time.sleep(1)
    
    def process_events(self, events):
        """处理收集到的事件"""
        # 防止索引正在进行时尝试更新
        if index_status["in_progress"]:
            logger.info("索引正在进行中，暂不处理文件变动")
            return
        
        # 收集所有需要添加和删除的文件
        files_to_update = set()
        files_to_remove = set()
        
        for event in events:
            if isinstance(event, FileCreatedEvent) or isinstance(event, FileModifiedEvent):
                # 使用is_supported_file检查文件是否应该被处理
                if is_supported_file(event.src_path):
                    # 检查文件大小
                    try:
                        is_too_large, file_size_mb = is_file_too_large(event.src_path)
                        file_name = os.path.basename(event.src_path)
                        # 对于特别大的文件，跳过处理
                        if is_too_large:
                            logger.warning(f"监控到的文件 '{file_name}' 过大 ({file_size_mb:.2f}MB > {MAX_FILE_SIZE_MB}MB)，已跳过")
                            continue
                        # 添加到更新列表
                        files_to_update.add(event.src_path)
                    except Exception as e:
                        logger.error(f"检查文件大小时出错: {str(e)}")
                        # 出错时还是添加到更新列表，让后续处理决定是否跳过
                        files_to_update.add(event.src_path)
            elif isinstance(event, FileDeletedEvent):
                # 对于删除事件，检查文件名是否应该被忽略
                file_name = os.path.basename(event.src_path)
                if (file_name.startswith('.') or file_name in ['.DS_Store', 'Thumbs.db', 'desktop.ini'] or 
                    file_name.endswith('~') or file_name.startswith('~$')):
                    continue  # 忽略系统隐藏文件的删除事件
                    
                files_to_remove.add(event.src_path)
            elif isinstance(event, FileMovedEvent):
                # 移动事件，首先检查源文件和目标文件是否应该被忽略
                file_name_src = os.path.basename(event.src_path)
                file_name_dest = os.path.basename(event.dest_path)
                
                # 检查是否为隐藏或系统文件
                if (file_name_src.startswith('.') or file_name_src in ['.DS_Store', 'Thumbs.db', 'desktop.ini'] or 
                    file_name_src.endswith('~') or file_name_src.startswith('~$')):
                    continue  # 忽略系统隐藏文件
                    
                # 忽略移动到隐藏文件的情况
                if (file_name_dest.startswith('.') or file_name_dest in ['.DS_Store', 'Thumbs.db', 'desktop.ini'] or 
                    file_name_dest.endswith('~') or file_name_dest.startswith('~$')):
                    continue  # 忽略系统隐藏文件
                    
                # 源文件是支持的类型就移除
                if is_supported_file(event.src_path):
                    files_to_remove.add(event.src_path)
                # 目标文件是支持的类型就添加
                if is_supported_file(event.dest_path):
                    # 检查目标文件大小
                    try:
                        is_too_large, file_size_mb = is_file_too_large(event.dest_path)
                        file_name = os.path.basename(event.dest_path)
                        # 对于特别大的文件，跳过处理
                        if is_too_large:
                            logger.warning(f"监控到的文件 '{file_name}' 过大 ({file_size_mb:.2f}MB > {MAX_FILE_SIZE_MB}MB)，已跳过")
                            continue
                        # 添加到更新列表
                        files_to_update.add(event.dest_path)
                    except Exception as e:
                        logger.error(f"检查文件大小时出错: {str(e)}")
                        # 出错时还是添加到更新列表，让后续处理决定是否跳过
                        files_to_update.add(event.dest_path)
        
        # 如果没有变化，则无需更新
        if not files_to_update and not files_to_remove:
            logger.info("没有需要更新的文件")
            return
        
        logger.info(f"开始处理文件变动: {len(files_to_update)} 个文件需更新, {len(files_to_remove)} 个文件需删除")
        
        try:
            # 检查向量数据库是否存在
            if not os.path.exists(self.db_path) or not os.path.exists(os.path.join(self.db_path, "index.faiss")):
                logger.warning(f"向量数据库不存在，需要完整重建索引: {self.folder_path}")
                # 触发完整索引重建
                background_tasks = BackgroundTasks()
                background_tasks.add_task(index_folder, self.folder_path)
                return
            
            # 加载现有向量库
            embedding_model = DashScopeEmbeddings(
                model=EMBEDDING_MODEL_NAME,
                dashscope_api_key=DASHSCOPE_API_KEY
            )
            
            # 加载现有数据库
            db = FAISS.load_local(self.db_path, embedding_model, allow_dangerous_deserialization=True)
            
            # 处理文件更新（添加和修改）
            if files_to_update:
                all_docs = []
                for file_path in files_to_update:
                    try:
                        logger.info(f"处理更新文件: {file_path}")
                        # 加载文档
                        file_docs = load_document(file_path)
                        if file_docs:
                            # 移除该文件的现有向量（如果存在）
                            try:
                                # 尝试使用filter删除
                                docs_to_keep = []
                                search_results = db.similarity_search_with_score("", k=1000)
                                for doc, _ in search_results:
                                    if doc.metadata.get("source") != file_path:
                                        docs_to_keep.append(doc)
                                
                                # 如果找到了需要保留的文档，重建索引
                                if docs_to_keep:
                                    logger.info(f"重建索引以删除文件 {file_path} 的文档")
                                    # 创建新的向量库
                                    new_db = FAISS.from_documents(docs_to_keep, embedding_model)
                                    # 保存新库到临时目录
                                    temp_path = f"{self.db_path}_temp"
                                    if os.path.exists(temp_path):
                                        shutil.rmtree(temp_path)
                                    os.makedirs(temp_path, exist_ok=True)
                                    new_db.save_local(temp_path)
                                    # 加载新库
                                    db = FAISS.load_local(temp_path, embedding_model, allow_dangerous_deserialization=True)
                                    # 删除临时库
                                    shutil.rmtree(temp_path)
                            except Exception as e:
                                logger.warning(f"删除文件 {file_path} 现有向量时出错: {str(e)}")
                                logger.info(f"继续处理更新，将添加文件的新内容")
                            
                            # 分割文档
                            text_splitter = CharacterTextSplitter(chunk_size=800, chunk_overlap=150)
                            split_docs = text_splitter.split_documents(file_docs)
                            all_docs.extend(split_docs)
                            logger.info(f"文件 {file_path} 更新成功")
                    except Exception as e:
                        logger.error(f"处理文件 {file_path} 更新失败: {str(e)}")
                
                # 将所有新文档添加到向量库
                if all_docs:
                    logger.info(f"添加 {len(all_docs)} 个文档块到向量库")
                    db.add_documents(all_docs)
            
            # 处理文件删除
            if files_to_remove:
                for file_path in files_to_remove:
                    try:
                        logger.info(f"从向量库中删除文件: {file_path}")
                        
                        try:
                            # 尝试使用filter重建索引，删除特定文件的文档
                            docs_to_keep = []
                            search_results = db.similarity_search_with_score("", k=1000)
                            for doc, _ in search_results:
                                if doc.metadata.get("source") != file_path:
                                    docs_to_keep.append(doc)
                            
                            # 如果找到了需要保留的文档，重建索引
                            if docs_to_keep:
                                logger.info(f"重建索引以删除文件 {file_path} 的文档")
                                # 创建新的向量库
                                new_db = FAISS.from_documents(docs_to_keep, embedding_model)
                                # 保存新库到临时目录
                                temp_path = f"{self.db_path}_temp"
                                if os.path.exists(temp_path):
                                    shutil.rmtree(temp_path)
                                os.makedirs(temp_path, exist_ok=True)
                                new_db.save_local(temp_path)
                                # 加载新库
                                db = FAISS.load_local(temp_path, embedding_model, allow_dangerous_deserialization=True)
                                # 删除临时库
                                shutil.rmtree(temp_path)
                                logger.info(f"文件 {file_path} 已从向量库中移除")
                            else:
                                logger.warning(f"未找到与文件 {file_path} 相关的文档")
                        except Exception as e:
                            logger.error(f"删除文件 {file_path} 的向量失败: {str(e)}")
                    except Exception as e:
                        logger.error(f"处理文件 {file_path} 删除操作失败: {str(e)}")
            
            # 保存更新后的向量库
            logger.info("保存更新后的向量库")
            db.save_local(self.db_path)
            logger.info("向量库更新完成")
            
        except Exception as e:
            error_msg = str(e)
            logger.error(f"处理文件变动时出错: {error_msg}")
            logger.debug(f"错误详情: {traceback.format_exc()}")

# 保存监控配置到文件
def save_monitoring_config():
    """保存当前的监控配置到文件"""
    try:
        config_dir = os.path.join(os.path.dirname(os.path.dirname(os.path.abspath(__file__))), "vector_store")
        os.makedirs(config_dir, exist_ok=True)
        config_file = os.path.join(config_dir, "monitoring_config.json")
        
        # 准备配置数据
        config = {}
        for folder in active_observers.keys():
            config[folder] = {"monitoring": True}
        
        # 保存配置
        with open(config_file, 'w', encoding='utf-8') as f:
            json.dump(config, f, ensure_ascii=False, indent=2)
            
        logger.info(f"已保存监控配置到 {config_file}")
    except Exception as e:
        logger.error(f"保存监控配置时出错: {str(e)}")

# 启动文件监控
def start_file_monitoring(folder_path):
    """为指定文件夹启动文件监控"""
    folder_path = normalize_path(folder_path)
    
    # 如果已经在监控，先停止
    stop_file_monitoring(folder_path)
    
    try:
        logger.info(f"启动文件监控: {folder_path}")
        event_handler = FileIndexHandler(folder_path)
        observer = Observer()
        observer.schedule(event_handler, folder_path, recursive=True)
        observer.start()
        
        # 保存到全局字典
        active_observers[folder_path] = observer
        file_handlers[folder_path] = event_handler
        
        # 保存更新后的监控配置
        save_monitoring_config()
        
        logger.info(f"文件监控已启动: {folder_path}")
        return True
    except Exception as e:
        logger.error(f"启动文件监控失败: {str(e)}")
        return False

# 停止文件监控
def stop_file_monitoring(folder_path):
    """停止指定文件夹的监控"""
    folder_path = normalize_path(folder_path)
    
    if folder_path in active_observers:
        try:
            logger.info(f"停止文件监控: {folder_path}")
            observer = active_observers[folder_path]
            observer.stop()
            observer.join()
            
            # 从全局字典中移除
            del active_observers[folder_path]
            if folder_path in file_handlers:
                del file_handlers[folder_path]
            
            # 保存更新后的监控配置
            save_monitoring_config()
                
            logger.info(f"文件监控已停止: {folder_path}")
            return True
        except Exception as e:
            logger.error(f"停止文件监控失败: {str(e)}")
            return False
    return False

# 停止所有文件监控
def stop_all_monitoring():
    """停止所有活跃的文件监控"""
    folders = list(active_observers.keys())
    for folder in folders:
        stop_file_monitoring(folder)
    return True

# 自定义Docx加载器
class CustomDocxLoader(BaseLoader):
    def __init__(self, file_path):
        self.file_path = file_path
    
    def load(self):
        try:
            # 检查文件后缀
            _, ext = os.path.splitext(self.file_path.lower())
            
            # 对于.doc文件（非.docx），直接返回错误
            if ext == '.doc':
                logger.error(f"不支持旧版Word格式(.doc)文件: {os.path.basename(self.file_path)}")
                logger.info("请将文件转换为.docx格式后再试")
                return []
            
            # 只处理.docx文件
            import docx
            doc = docx.Document(self.file_path)
            
            # 提取文本
            paragraphs = []
            for para in doc.paragraphs:
                if para.text and para.text.strip():
                    paragraphs.append(para.text)
            
            # 提取表格内容
            for table in doc.tables:
                table_text = []
                for row in table.rows:
                    row_text = []
                    for cell in row.cells:
                        if cell.text and cell.text.strip():
                            row_text.append(cell.text.strip())
                    if row_text:
                        table_text.append(" | ".join(row_text))
                if table_text:
                    paragraphs.append("\n".join(table_text))
            
            text = "\n\n".join(paragraphs)
            
            # 如果没有提取到任何文本，返回空列表
            if not text or not text.strip():
                logger.warning(f"从文件 {os.path.basename(self.file_path)} 中未提取到任何文本")
                return []
                
            metadata = {"source": self.file_path}
            return [Document(page_content=text, metadata=metadata)]
        except Exception as e:
            logger.error(f"读取Docx文件时出错: {str(e)}")
            return []

# 自定义PPT加载器
class CustomPptxLoader(BaseLoader):
    def __init__(self, file_path):
        self.file_path = file_path
    
    def load(self):
        try:
            from pptx import Presentation
            prs = Presentation(self.file_path)
            
            text_parts = []
            for slide in prs.slides:
                slide_text = []
                for shape in slide.shapes:
                    if hasattr(shape, "text") and shape.text:
                        slide_text.append(shape.text)
                if slide_text:
                    text_parts.append("\n".join(slide_text))
            
            full_text = "\n\n".join(text_parts)
            metadata = {"source": self.file_path}
            return [Document(page_content=full_text, metadata=metadata)]
        except Exception as e:
            logger.error(f"读取PPT文件时出错: {str(e)}")
            return []

# 自定义Excel加载器
class CustomExcelLoader(BaseLoader):
    def __init__(self, file_path):
        self.file_path = file_path
    
    def load(self):
        try:
            import pandas as pd
            
            # 读取所有工作表
            xls = pd.ExcelFile(self.file_path)
            all_sheet_texts = []
            
            for sheet_name in xls.sheet_names:
                try:
                    # 读取工作表
                    df = pd.read_excel(self.file_path, sheet_name=sheet_name)
                    
                    # 将NaN值替换为空字符串
                    df = df.fillna("")
                    
                    # 转换为字符串格式
                    sheet_text = f"工作表: {sheet_name}\n\n"
                    # 添加列名
                    sheet_text += "列: " + ", ".join(str(col) for col in df.columns) + "\n\n"
                    
                    # 添加每行数据
                    for idx, row in df.iterrows():
                        row_text = " | ".join(str(val) for val in row.values)
                        sheet_text += f"行 {idx+1}: {row_text}\n"
                    
                    all_sheet_texts.append(sheet_text)
                except Exception as e:
                    logger.error(f"读取Excel工作表 {sheet_name} 时出错: {str(e)}")
            
            # 如果成功解析了任何工作表
            if all_sheet_texts:
                # 组合所有工作表的文本
                full_text = "\n\n".join(all_sheet_texts)
                return [Document(page_content=full_text, metadata={"source": self.file_path})]
            return []
        except Exception as e:
            logger.error(f"读取Excel文件时出错: {str(e)}")
            return []

# 导入dashscope模块并配置其日志级别
import dashscope
import dashscope.embeddings
dashscope.api_key = DASHSCOPE_API_KEY  # 使用通义千问API密钥
# 配置dashscope客户端日志级别 (如果有这样的API)
# 注意: dashscope可能没有直接的日志控制API
# 但我们已经在下面设置了关联模块的日志级别

# 加载环境变量

# 初始化FastAPI应用
app = FastAPI(title="FindMe API")

# 添加错误处理
@app.exception_handler(Exception)
async def global_exception_handler(request: Request, exc: Exception):
    error_msg = f"全局错误: {str(exc)}"
    trace = traceback.format_exc()
    logger.error(f"{error_msg}\n{trace}")
    return JSONResponse(
        status_code=500,
        content={"success": False, "message": str(exc), "detail": trace},
    )

# 添加中间件处理请求
@app.middleware("http")
async def log_requests(request: Request, call_next):
    logger.info(f"请求: {request.method} {request.url.path}")
    try:
        # 处理请求
        response = await call_next(request)
        logger.info(f"响应: {response.status_code}")
        return response
    except Exception as e:
        logger.error(f"处理请求时出错: {e}")
        return JSONResponse(
            status_code=500,
            content={"success": False, "message": f"处理请求时出错: {str(e)}"}
        )

# 添加CORS中间件
app.add_middleware(
    CORSMiddleware,
    allow_origins=["*"],
    allow_credentials=True,
    allow_methods=["*"],
    allow_headers=["*"],
)

# 请求模型定义
class FolderRequest(BaseModel):
    folder: str

class SearchRequest(BaseModel):
    query: str
    folder: str

class FileRequest(BaseModel):
    file_path: str
    
class ConfigRequest(BaseModel):
    max_text_length: Optional[int] = None
    max_chunk_count: Optional[int] = None
    max_file_size_mb: Optional[int] = None
    embedding_model: Optional[str] = None  # 添加嵌入模型选择

# 工具函数：规范化路径
def normalize_path(path: str) -> str:
    """规范化路径，处理特殊字符"""
    # 确保路径使用正确的分隔符
    normalized = os.path.normpath(path).replace('\\', '/')
    return normalized

# 工具函数：获取向量数据库路径
def get_db_path(folder: str) -> str:
    """根据文件夹路径生成向量数据库路径"""
    # 使用MD5哈希值作为数据库文件夹名，更稳定
    import hashlib
    # 确保输入是 str 类型
    if not isinstance(folder, str):
        folder = str(folder)
    folder = normalize_path(folder)
    folder_hash = hashlib.md5(folder.encode('utf-8')).hexdigest()
    db_path = os.path.join(os.path.dirname(os.path.dirname(os.path.abspath(__file__))), "vector_store", folder_hash)
    return db_path

# 工具函数：检查文件类型是否支持
def is_supported_file(file_path: str) -> bool:
    """检查文件类型是否支持索引"""
    supported_extensions = ['.txt', '.pdf', '.docx', '.pptx', '.xlsx', '.xls', '.csv']  # 添加了.csv支持
    try:
        # 获取文件名和扩展名
        file_name = os.path.basename(file_path)
        _, ext = os.path.splitext(file_path.lower())
        
        # 过滤系统隐藏文件
        # 1. 过滤.DS_Store文件（macOS）
        if file_name == '.DS_Store':
            return False
        # 2. 过滤以.开头的隐藏文件
        if file_name.startswith('.'):
            return False
        # 3. 过滤Windows系统文件
        if file_name in ['Thumbs.db', 'desktop.ini']:
            return False
        # 4. 过滤临时文件
        if file_name.endswith('~') or file_name.startswith('~$') or ext in ['.tmp', '.temp']:
            return False
            
        # 检查是否是支持的扩展名和有效文件
        return ext in supported_extensions and os.path.isfile(file_path)
    except Exception as e:
        logger.error(f"检查文件类型时出错: {file_path}, {e}")
        return False

# 工具函数：检查文件大小是否超过限制
def is_file_too_large(file_path: str) -> (bool, float):
    """检查文件大小是否超过限制，返回(是否超过限制, 文件大小MB)"""
    try:
        file_size = os.path.getsize(file_path)
        file_size_mb = file_size / (1024 * 1024)
        return file_size_mb > MAX_FILE_SIZE_MB, file_size_mb
    except Exception as e:
        logger.error(f"检查文件大小时出错: {file_path}, {e}")
        return False, 0

# 工具函数：加载文档
def load_document(file_path: str) -> List:
    """根据文件类型加载文档，并应用字数限制"""
    try:
        _, ext = os.path.splitext(file_path.lower())
        
        # 加载文档
        docs = []
        if ext == '.txt':
            docs = [TextLoader(file_path).load()[0]]
        elif ext == '.pdf':
            try:
                docs = PyPDFLoader(file_path).load()
            except Exception as pdf_error:
                if "Odd-length string" in str(pdf_error):
                    logger.error(f"加载PDF文件出错 {os.path.basename(file_path)}: Odd-length string")
                    # 尝试使用其他方式提取PDF内容
                    try:
                        import fitz  # PyMuPDF
                        text_content = ""
                        with fitz.open(file_path) as pdf:
                            for page_num in range(len(pdf)):
                                text_content += pdf[page_num].get_text()
                        
                        if text_content.strip():
                            metadata = {"source": file_path}
                            docs = [Document(page_content=text_content, metadata=metadata)]
                            logger.info(f"使用PyMuPDF成功加载PDF: {os.path.basename(file_path)}")
                        else:
                            logger.error(f"使用PyMuPDF提取的PDF内容为空: {os.path.basename(file_path)}")
                    except ImportError:
                        logger.error("PyMuPDF (fitz) 库未安装，无法使用备选方法解析PDF")
                    except Exception as muPdf_error:
                        logger.error(f"使用PyMuPDF加载PDF文件失败: {str(muPdf_error)}")
                else:
                    # 其他PDF解析错误
                    raise pdf_error
        elif ext == '.docx':
            docs = CustomDocxLoader(file_path).load()
        elif ext == '.pptx':
            docs = CustomPptxLoader(file_path).load()
        elif ext in ['.xlsx', '.xls']:
            docs = CustomExcelLoader(file_path).load()
        elif ext == '.csv':
            # 处理CSV文件
            try:
                import csv
                import io
                
                # 尝试检测编码
                encoding = 'utf-8'  # 默认编码
                try:
                    import chardet
                    with open(file_path, 'rb') as f:
                        result = chardet.detect(f.read())
                        if result['confidence'] > 0.7:  # 仅当置信度高于0.7时使用检测结果
                            encoding = result['encoding']
                    logger.info(f"检测到CSV文件编码: {encoding}")
                except ImportError:
                    logger.warning("chardet库未安装，使用默认utf-8编码")
                except Exception as e:
                    logger.warning(f"检测CSV文件编码失败: {str(e)}，使用默认utf-8编码")
                
                # 读取CSV文件
                with open(file_path, 'r', encoding=encoding, errors='replace') as f:
                    csv_reader = csv.reader(f)
                    rows = list(csv_reader)
                    
                    if not rows:
                        logger.warning(f"CSV文件为空: {os.path.basename(file_path)}")
                        return []
                    
                    # 提取表头（第一行）
                    headers = rows[0] if rows else []
                    
                    # 构建文本内容
                    text_content = ""
                    
                    # 添加表头信息
                    if headers:
                        text_content += f"表头: {' | '.join(headers)}\n\n"
                    
                    # 添加每行数据
                    for i, row in enumerate(rows[1:], 1):  # 跳过表头，从第二行开始
                        # 将行数据合并为文本
                        row_values = [str(val).strip() for val in row]
                        row_text = " | ".join(row_values)
                        text_content += f"行 {i}: {row_text}\n"
                    
                    metadata = {"source": file_path}
                    docs = [Document(page_content=text_content, metadata=metadata)]
                    logger.info(f"成功加载CSV文件: {os.path.basename(file_path)}")
            except Exception as e:
                logger.error(f"加载CSV文件失败 {os.path.basename(file_path)}: {str(e)}")
                return []
        
        # 应用文本长度限制
        if not docs:
            return []
        
        # 获取文件名
        file_name = os.path.basename(file_path)
        
        # 处理文档，添加文件名到内容开头
        limited_docs = []
        for doc in docs:
            # 原始内容
            content = doc.page_content
            
            # 添加文件名到内容开头
            enhanced_content = f"文件: {file_name}\n\n{content}"
            
            # 检查长度限制
            if len(enhanced_content) > MAX_TEXT_LENGTH:
                logger.warning(f"文件 '{file_name}' 内容过长，已截断至 {MAX_TEXT_LENGTH} 字符")
                # 创建一个新的文档对象，使用截断的内容
                limited_content = enhanced_content[:MAX_TEXT_LENGTH] + "\n[内容过长，已截断...]"
                limited_docs.append(Document(page_content=limited_content, metadata=doc.metadata))
            else:
                limited_docs.append(Document(page_content=enhanced_content, metadata=doc.metadata))
        
        return limited_docs
    except Exception as e:
        logger.error(f"加载文件出错 {file_path}: {str(e)}")
        return []

# 索引文件夹中的文档
def index_folder(folder: str):
    """索引文件夹中的所有支持的文档"""
    global index_status
    
    try:
        logger.info(f"开始索引文件夹: {folder}")
        index_status["in_progress"] = True
        index_status["progress"] = 0
        index_status["status"] = "准备索引..."
        index_status["error"] = None
        index_status["completed"] = False
        index_status["file_stats"] = {
            "success_count": 0,
            "failure_count": 0,
            "skipped_count": 0,
            "total_count": 0
        }
        index_status["success_files"] = []
        index_status["failed_files"] = []
        index_status["skipped_files"] = []
        
        # 规范化并检查路径
        folder = normalize_path(folder)
        if not os.path.exists(folder):
            raise ValueError(f"文件夹不存在: {folder}")
        if not os.path.isdir(folder):
            raise ValueError(f"路径不是文件夹: {folder}")
        
        # 获取数据库路径
        db_path = get_db_path(folder)
        logger.info(f"数据库路径: {db_path}")
        
        # 初始化变量
        indexed_files = {}
        db = None
        
        # 如果数据库已存在，检查现有索引并进行增量更新
        if os.path.exists(db_path) and os.path.exists(os.path.join(db_path, "index.faiss")):
            logger.info(f"发现现有向量数据库，将进行增量更新: {db_path}")
            index_status["status"] = "发现现有索引，准备增量更新..."
            
            # 加载现有数据库
            try:
                embedding_model = DashScopeEmbeddings(
                    model=EMBEDDING_MODEL_NAME,
                    dashscope_api_key=DASHSCOPE_API_KEY
                )
                db = FAISS.load_local(db_path, embedding_model, allow_dangerous_deserialization=True)
                
                # 获取已索引文件列表
                # 查询所有文档，获取源文件路径
                search_results = db.similarity_search_with_score("", k=10000)  # 尝试获取所有文档
                for doc, _ in search_results:
                    source = doc.metadata.get("source", "")
                    if source and os.path.exists(source):
                        # 获取文件的最后修改时间
                        mtime = os.path.getmtime(source)
                        if source not in indexed_files or mtime > indexed_files[source]["mtime"]:
                            indexed_files[source] = {
                                "mtime": mtime,
                                "indexed": True
                            }
                
                logger.info(f"从现有索引中找到 {len(indexed_files)} 个已索引文件")
                index_status["status"] = f"从现有索引中找到 {len(indexed_files)} 个已索引文件，准备增量更新..."
            except Exception as e:
                logger.error(f"加载现有索引出错: {str(e)}")
                logger.info("将重新创建索引数据库")
                # 如果加载失败，删除现有索引并重建
                if os.path.exists(db_path):
                    shutil.rmtree(db_path)
                os.makedirs(db_path, exist_ok=True)
                indexed_files = {}
                db = None
        else:
            # 确保数据库目录存在
            os.makedirs(db_path, exist_ok=True)
        
        # 扫描文件夹中的所有文件 - 不再过滤格式，收集所有文件
        all_files = []
        unsupported_files = []
        
        for root, _, files in os.walk(folder):
            for file in files:
                try:
                    file_path = os.path.join(root, file)
                    # 只检查是否为文件，不再使用is_supported_file过滤
                    if os.path.isfile(file_path):
                        # 检查是否支持该格式
                        if is_supported_file(file_path):
                            all_files.append(file_path)
                        else:
                            # 检查是否是系统文件（如.DS_Store）
                            file_name = os.path.basename(file_path)
                            # 如果不是系统隐藏文件，则添加到不支持文件列表
                            if not (file_name == '.DS_Store' or 
                                   file_name.startswith('.') or 
                                   file_name in ['Thumbs.db', 'desktop.ini'] or 
                                   file_name.endswith('~') or 
                                   file_name.startswith('~$')):
                                unsupported_files.append(file_path)
                except Exception as e:
                    logger.error(f"扫描文件时出错: {file}, {e}")
        
        # 计算总文件数，包括不支持的格式
        total_files = len(all_files) + len(unsupported_files)
        
        if total_files == 0:
            index_status["status"] = "没有找到任何文件"
            index_status["completed"] = True
            index_status["in_progress"] = False
            logger.warning(f"文件夹 {folder} 中没有找到任何文件")
            return
        
        index_status["status"] = f"找到 {total_files} 个文件（其中 {len(unsupported_files)} 个格式不支持），开始加载..."
        index_status["file_stats"]["total_count"] = total_files
        logger.info(f"找到 {total_files} 个文件（其中 {len(unsupported_files)} 个格式不支持），开始加载...")
        
        # 加载所有文档
        docs = []
        success_count = 0
        failed_files = []
        skipped_files = []
        embedding_model = None  # 初始化嵌入模型变量
        
        # 先处理不支持的文件格式，将它们标记为失败
        for file_path in unsupported_files:
            file_name = os.path.basename(file_path)
            _, ext = os.path.splitext(file_path.lower())
            logger.warning(f"不支持的文件格式 '{ext}': {file_name}")
            failed_files.append(f"{file_name} (不支持的格式: {ext})")
            index_status["failed_files"].append({
                "name": file_name,
                "path": file_path,
                "reason": f"不支持的格式: {ext}"
            })
            index_status["file_stats"]["failure_count"] += 1
        
        # 处理支持的文件格式
        for i, file_path in enumerate(all_files):
            # 计算进度，考虑到已处理的不支持文件格式
            processed_count = len(unsupported_files) + i
            index_status["progress"] = int((processed_count / total_files) * 50)  # 前50%进度用于加载文件
            file_name = os.path.basename(file_path)
            index_status["status"] = f"加载文件 ({processed_count+1}/{total_files}): {file_name}"
            
            # 减少日志输出，只在每10个文件或最后一个文件时记录日志
            if i % 10 == 0 or i == len(all_files) - 1:
                logger.info(f"加载文件 {processed_count+1}/{total_files}: {file_name}")
            
            try:
                # 检查文件是否已存在于索引中并且没有更改
                current_mtime = os.path.getmtime(file_path)
                if file_path in indexed_files and current_mtime <= indexed_files[file_path]["mtime"]:
                    logger.info(f"文件 '{file_name}' 已存在于索引中且未更改，跳过处理")
                    success_count += 1
                    index_status["success_files"].append({
                        "name": file_name,
                        "path": file_path,
                        "skipped": True,
                        "reason": "已索引且未更改"
                    })
                    index_status["file_stats"]["success_count"] += 1
                    continue
                
                # 检查文件大小
                is_too_large, file_size_mb = is_file_too_large(file_path)
                # 对于特别大的文件，跳过处理
                if is_too_large:
                    logger.warning(f"文件 '{file_name}' 过大 ({file_size_mb:.2f}MB > {MAX_FILE_SIZE_MB}MB)，已跳过")
                    skipped_files.append(f"{file_name} (过大: {file_size_mb:.2f}MB)")
                    index_status["skipped_files"].append({
                        "name": file_name,
                        "path": file_path,
                        "reason": f"文件过大: {file_size_mb:.2f}MB"
                    })
                    index_status["file_stats"]["skipped_count"] += 1
                    continue
                
                file_docs = load_document(file_path)
                if file_docs:
                    # 限制单个文件的块数量
                    if len(file_docs) > MAX_CHUNK_COUNT:
                        logger.warning(f"文件 '{file_name}' 生成的块数 ({len(file_docs)}) 超过限制 ({MAX_CHUNK_COUNT})，已截断")
                        file_docs = file_docs[:MAX_CHUNK_COUNT]
                        
                    docs.extend(file_docs)
                    success_count += 1
                    index_status["success_files"].append({
                        "name": file_name,
                        "path": file_path
                    })
                    index_status["file_stats"]["success_count"] += 1
                    
                    # 更新已索引文件记录
                    indexed_files[file_path] = {
                        "mtime": current_mtime,
                        "indexed": True
                    }
                    
                    # 流水线处理 - 每处理一定数量的文件就进行一次向量化
                    # 设置阈值，每积累100个文档块或处理了20个文件就进行一次向量化
                    if len(docs) >= 100 or (i > 0 and i % 20 == 0):
                        # 文本切片
                        partial_progress = int((processed_count / total_files) * 50)
                        index_status["status"] = f"已处理 {success_count} 个文件，进行部分向量化..."
                        index_status["progress"] = 50 + partial_progress // 2  # 进度表现为50%-75%之间
                        
                        # 分块处理
                        text_splitter = CharacterTextSplitter(chunk_size=800, chunk_overlap=150)
                        split_docs = text_splitter.split_documents(docs)
                        
                        # 检查并修剪文本块
                        valid_split_docs = []
                        for doc in split_docs:
                            if len(doc.page_content) > MAX_TEXT_BLOCK_SIZE:
                                # 截断超长内容
                                truncated_content = doc.page_content[:(MAX_TEXT_BLOCK_SIZE-8)] + "..."
                                valid_split_docs.append(Document(page_content=truncated_content, metadata=doc.metadata))
                            else:
                                valid_split_docs.append(doc)
                        
                        # 创建嵌入模型（如果还没创建）
                        if embedding_model is None:
                            logger.info("初始化嵌入模型...")
                            embedding_model = DashScopeEmbeddings(
                                model=EMBEDDING_MODEL_NAME,
                                dashscope_api_key=DASHSCOPE_API_KEY
                            )
                        
                        # 向量化处理
                        logger.info(f"中间向量化处理 {len(valid_split_docs)} 个文本块...")
                        try:
                            # 分批向量化
                            batch_size = MAX_BATCH_ROWS
                            for j in range(0, len(valid_split_docs), batch_size):
                                end_j = min(j + batch_size, len(valid_split_docs))
                                batch_docs = valid_split_docs[j:end_j]
                                try:
                                    if db is None:
                                        # 第一批，创建新的向量数据库
                                        db = FAISS.from_documents(batch_docs, embedding_model)
                                    else:
                                        # 向现有数据库添加文档
                                        db.add_documents(batch_docs)
                                    logger.info(f"成功向量化 {end_j-j} 个文本块")
                                except Exception as batch_error:
                                    logger.error(f"批量向量化出错: {str(batch_error)}")
                                    # 尝试减小批量
                                    retry_size = len(batch_docs) // 2
                                    if retry_size > 0:
                                        logger.info(f"尝试以较小批量 ({retry_size}) 重试...")
                                        for k in range(0, len(batch_docs), retry_size):
                                            retry_docs = batch_docs[k:k+retry_size]
                                            try:
                                                if db is None:
                                                    db = FAISS.from_documents(retry_docs, embedding_model)
                                                else:
                                                    db.add_documents(retry_docs)
                                            except Exception:
                                                # 继续处理下一批
                                                continue
                        except Exception as e:
                            logger.error(f"中间向量化处理失败: {str(e)}")
                        
                        # 保存中间结果
                        if db is not None:
                            try:
                                logger.info("保存中间向量化结果...")
                                db.save_local(db_path)
                            except Exception as save_error:
                                logger.error(f"保存中间结果失败: {str(save_error)}")
                        
                        # 清空处理队列，开始下一批文件
                        docs = []
                else:
                    logger.error(f"文件解析结果为空: {file_name}")
                    failed_files.append(file_name)
                    index_status["failed_files"].append({
                        "name": file_name,
                        "path": file_path,
                        "reason": "解析结果为空"
                    })
                    index_status["file_stats"]["failure_count"] += 1
            except Exception as e:
                logger.error(f"加载文件失败 {file_name}: {str(e)}")
                failed_files.append(file_name)
                index_status["failed_files"].append({
                    "name": file_name,
                    "path": file_path,
                    "reason": f"加载失败: {str(e)}"
                })
                index_status["file_stats"]["failure_count"] += 1
        
        # 统计并显示成功和失败的文件
        failure_count = len(failed_files)
        skipped_count = len(skipped_files)
        
        # 更新索引状态的文件统计
        index_status["file_stats"]["success_count"] = success_count
        index_status["file_stats"]["failure_count"] = failure_count
        index_status["file_stats"]["skipped_count"] = skipped_count
        
        logger.info(f"文件处理完成: 成功 {success_count} 个, 失败 {failure_count} 个, 跳过 {skipped_count} 个")
        if failed_files:
            logger.info(f"失败的文件: {', '.join(failed_files[:10])}" + 
                       (f" 等 {len(failed_files)} 个文件" if len(failed_files) > 10 else ""))
        if skipped_files:
            logger.info(f"跳过的文件: {', '.join(skipped_files[:10])}" + 
                       (f" 等 {len(skipped_files)} 个文件" if len(skipped_files) > 10 else ""))
        
        # 处理队列中剩余的文档（最后一批）
        if docs:
            # 文本切片
            index_status["status"] = f"处理最后一批文件..."
            index_status["progress"] = 75
            logger.info("处理剩余文档...")
            
            # 分块处理
            text_splitter = CharacterTextSplitter(chunk_size=800, chunk_overlap=150)
            split_docs = text_splitter.split_documents(docs)
            
            # 检查并修剪文本块
            valid_split_docs = []
            for doc in split_docs:
                if len(doc.page_content) > MAX_TEXT_BLOCK_SIZE:
                    truncated_content = doc.page_content[:(MAX_TEXT_BLOCK_SIZE-8)] + "..."
                    valid_split_docs.append(Document(page_content=truncated_content, metadata=doc.metadata))
                else:
                    valid_split_docs.append(doc)
            
            # 截断过多的文本块，避免请求数量过大
            if len(valid_split_docs) > MAX_CHUNK_COUNT:
                logger.warning(f"文本块数量超出限制，从 {len(valid_split_docs)} 截断至 {MAX_CHUNK_COUNT}")
                valid_split_docs = valid_split_docs[:MAX_CHUNK_COUNT]
            
            # 创建嵌入模型（如果还没有创建）
            if embedding_model is None:
                index_status["status"] = "加载嵌入模型..."
                index_status["progress"] = 80
                logger.info("加载通义千问嵌入模型...")
                
                embedding_model = DashScopeEmbeddings(
                    model=EMBEDDING_MODEL_NAME,
                    dashscope_api_key=DASHSCOPE_API_KEY
                )
            
            # 构建向量数据库
            index_status["status"] = "处理最后批次文本块..."
            index_status["progress"] = 85
            logger.info(f"对剩余 {len(valid_split_docs)} 个文本块创建向量索引...")
            
            # 分批处理文本块
            batch_size = MAX_BATCH_ROWS
            for i in range(0, len(valid_split_docs), batch_size):
                end_idx = min(i + batch_size, len(valid_split_docs))
                batch_docs = valid_split_docs[i:end_idx]
                
                try:
                    if db is None:
                        db = FAISS.from_documents(batch_docs, embedding_model)
                    else:
                        db.add_documents(batch_docs)
                except Exception as e:
                    logger.error(f"处理最后批次时出错: {str(e)}")
                    # 如果批处理失败，尝试减小批大小并重试
                    retry_batch_size = len(batch_docs) // 2
                    if retry_batch_size > 0:
                        for j in range(0, len(batch_docs), retry_batch_size):
                            try:
                                retry_docs = batch_docs[j:j+retry_batch_size]
                                if db is None:
                                    db = FAISS.from_documents(retry_docs, embedding_model)
                                else:
                                    db.add_documents(retry_docs)
                            except:
                                continue
        else:
            # 如果没有文档需要处理
            if not db:  # 如果也没有中间生成的数据库
                index_status["status"] = f"未能成功加载任何文件。尝试处理了 {total_files} 个文件，全部失败或跳过。"
                index_status["completed"] = True
                index_status["in_progress"] = False
                logger.warning("没有成功加载任何文件")
                return

        if db is None:
            raise Exception("所有批次处理均失败，无法创建向量数据库")

        # 保存向量数据库
        index_status["status"] = "保存向量数据库..."
        index_status["progress"] = 90
        logger.info("保存最终向量数据库...")

        db.save_local(db_path)

        # 索引完成
        index_status["status"] = f"索引完成！成功处理 {success_count} 个文件，失败 {failure_count} 个文件，跳过 {skipped_count} 个文件。"
        index_status["progress"] = 100
        index_status["completed"] = True
        logger.info(f"索引完成！成功处理 {success_count} 个文件，失败 {failure_count} 个文件，跳过 {skipped_count} 个文件。")
        
    except Exception as e:
        error_msg = str(e)
        trace = traceback.format_exc()
        logger.error(f"索引出错: {error_msg}")
        # 仅在DEBUG级别记录完整堆栈跟踪
        logger.debug(f"错误详情: {trace}")
        index_status["error"] = error_msg
        index_status["status"] = f"索引出错: {error_msg}"
    finally:
        index_status["in_progress"] = False

# API路由：检查数据库是否存在 - 改为POST请求
@app.post("/check-db")
async def check_db(folder_req: FolderRequest):
    folder = folder_req.folder
    logger.info(f"检查数据库是否存在: {folder}")
    try:
        # 规范化路径
        folder = normalize_path(folder)
        db_path = get_db_path(folder)
        exists = os.path.exists(db_path) and os.path.exists(os.path.join(db_path, "index.faiss"))
        logger.info(f"数据库路径: {db_path}, 存在: {exists}")
        return {"exists": exists}
    except Exception as e:
        logger.error(f"检查数据库出错: {str(e)}\n{traceback.format_exc()}")
        raise

# 保留GET方法兼容
@app.get("/check-db")
async def check_db_get(folder: str):
    logger.info(f"[GET] 检查数据库是否存在: {folder}")
    try:
        # 解码URL编码的路径
        folder = urllib.parse.unquote(folder)
        # 规范化路径
        folder = normalize_path(folder)
        db_path = get_db_path(folder)
        exists = os.path.exists(db_path) and os.path.exists(os.path.join(db_path, "index.faiss"))
        logger.info(f"数据库路径: {db_path}, 存在: {exists}")
        return {"exists": exists}
    except Exception as e:
        logger.error(f"检查数据库出错: {str(e)}\n{traceback.format_exc()}")
        raise

# API路由：开始索引
@app.post("/index")
async def start_index(folder_req: FolderRequest, background_tasks: BackgroundTasks):
    folder = folder_req.folder
    logger.info(f"收到索引请求: {folder}")
    
    try:
        # 规范化路径
        folder = normalize_path(folder)
        
        # 检查路径是否存在
        if not os.path.exists(folder):
            error = f"路径不存在: {folder}"
            logger.error(error)
            return {"success": False, "message": error}
        
        # 检查是否是文件夹
        if not os.path.isdir(folder):
            error = f"路径不是一个有效的文件夹: {folder}"
            logger.error(error)
            return {"success": False, "message": error}
        
        # 如果已经在索引中，则返回错误
        if index_status["in_progress"]:
            error = "已有索引任务在进行中，请等待完成后再试"
            logger.warning(error)
            return {"success": False, "message": error}
        
        # 检查是否已经有现有索引
        db_path = get_db_path(folder)
        has_existing_index = os.path.exists(db_path) and os.path.exists(os.path.join(db_path, "index.faiss"))
        if has_existing_index:
            logger.info(f"文件夹 {folder} 已存在索引，将进行增量更新")
        
        # 开始索引
        background_tasks.add_task(index_folder, folder)
        
        # 同时启动文件监控
        start_file_monitoring(folder)
        
        logger.info(f"索引任务已开始: {folder}")
        return {"success": True, "message": "开始增量索引更新..." if has_existing_index else "开始创建索引..."}
    except Exception as e:
        error_msg = str(e)
        logger.error(f"启动索引时出错: {error_msg}")
        return {"success": False, "message": f"启动索引时出错: {error_msg}"}

# API路由：获取索引进度
@app.get("/index-progress")
async def get_index_progress():
    return {
        "in_progress": index_status["in_progress"],
        "progress": index_status["progress"],
        "status": index_status["status"],
        "error": index_status["error"],
        "completed": index_status["completed"],
        "file_stats": index_status["file_stats"],
        "success_files": index_status["success_files"],
        "failed_files": index_status["failed_files"],
        "skipped_files": index_status["skipped_files"]
    }

# API路由：搜索
@app.post("/search")
async def search(search_req: SearchRequest):
    query = search_req.query
    folder = search_req.folder
    logger.info(f"搜索请求: '{query}', 文件夹: {folder}")
    
    try:
        # 规范化路径
        folder = normalize_path(folder)
        db_path = get_db_path(folder)
        
        # 检查数据库是否存在
        if not os.path.exists(db_path) or not os.path.exists(os.path.join(db_path, "index.faiss")):
            logger.error(f"向量数据库不存在: {db_path}")
            return {"success": False, "message": "向量数据库不存在，请先索引文件夹"}
        
        # 加载向量数据库
        embedding_model = DashScopeEmbeddings(
            model=EMBEDDING_MODEL_NAME,
            dashscope_api_key=DASHSCOPE_API_KEY
        )
        db = FAISS.load_local(db_path, embedding_model, allow_dangerous_deserialization=True)
        
        # 执行搜索 - 使用批处理模式，充分利用模型的行数限制
        # 默认情况下执行单次查询，因为查询只有一行文本
        try:
            # 标准查询模式
            docs_and_scores = db.similarity_search_with_score(query, k=MAX_BATCH_ROWS)
        except Exception as search_error:
            logger.warning(f"初始查询失败，尝试备用方法: {str(search_error)}")
            # 如果初始查询失败，使用更保守的k值
            try:
                docs_and_scores = db.similarity_search_with_score(query, k=10)
            except Exception as fallback_error:
                logger.error(f"备用查询也失败: {str(fallback_error)}")
                docs_and_scores = []
        
        # 按文件源去重，保留每个文件最相似的结果
        unique_sources = {}
        for doc, score in docs_and_scores:
            source = doc.metadata.get("source", "未知文件")
            float_score = float(score)
            
            # 如果文件还未记录，或当前结果比已记录的相似度更高
            if source not in unique_sources or float_score < unique_sources[source]["score"]:
                # 处理高亮
                content = doc.page_content
                highlighted_content = content
                
                # 分割查询词并找出最长的词组匹配
                query_terms = query.split()
                
                # 先尝试查找完整查询
                if query in content:
                    # 添加高亮标记 - 使用HTML标签作为高亮标记，前端可以解析它
                    highlighted_content = content.replace(query, f"<mark>{query}</mark>")
                else:
                    # 尝试查找较长的词组
                    for i in range(len(query_terms), 0, -1):
                        for j in range(len(query_terms) - i + 1):
                            phrase = " ".join(query_terms[j:j+i])
                            if phrase and len(phrase) > 1 and phrase in content:
                                # 避免重复替换已高亮的部分
                                highlighted_content = highlighted_content.replace(
                                    phrase, 
                                    f"<mark>{phrase}</mark>"
                                )
                
                unique_sources[source] = {
                    "content": content,
                    "highlighted_content": highlighted_content,
                    "source": source,
                    "score": float_score
                }
        
        # 将去重后的结果按相似度排序
        results = list(unique_sources.values())
        results.sort(key=lambda x: x["score"])
        
        # 限制返回数量
        results = results[:5]
        
        logger.info(f"查询完成，找到 {len(results)} 个结果")
        return {"success": True, "results": results}
    except Exception as e:
        error_msg = str(e)
        trace = traceback.format_exc()
        logger.error(f"搜索出错: {error_msg}")
        # 在错误情况下才记录完整堆栈跟踪
        logger.debug(f"错误详情: {trace}")
        return {"success": False, "message": f"搜索出错: {error_msg}"}

# API路由：打开文件
@app.post("/open-file")
async def open_file(file_req: FileRequest):
    file_path = file_req.file_path
    logger.info(f"请求打开文件: {file_path}")
    
    try:
        # 规范化路径
        file_path = normalize_path(file_path)
        
        if not os.path.exists(file_path):
            error = f"文件不存在: {file_path}"
            logger.error(error)
            return {"success": False, "message": error}
        
        # 根据操作系统使用不同的命令打开文件
        if sys.platform.startswith('darwin'):  # macOS
            subprocess.run(['open', file_path], check=True)
        elif sys.platform.startswith('win'):   # Windows
            os.startfile(file_path)
        else:  # Linux或其他
            subprocess.run(['xdg-open', file_path], check=True)
        
        logger.info(f"已成功打开文件: {file_path}")
        return {"success": True}
    except Exception as e:
        error_msg = str(e)
        trace = traceback.format_exc()
        logger.error(f"打开文件出错: {error_msg}")
        return {"success": False, "message": f"打开文件出错: {error_msg}"}

# 健康检查端点
@app.get("/health")
async def health_check():
    return {"status": "ok"}

# 调试端点，检查路径
@app.post("/debug/check-path")
async def debug_check_path(folder_req: FolderRequest):
    folder = folder_req.folder
    try:
        # 规范化路径
        normalized = normalize_path(folder)
        # 检查路径是否存在
        exists = os.path.exists(normalized)
        # 检查是否是文件夹
        is_dir = os.path.isdir(normalized) if exists else False
        # 获取数据库路径
        db_path = get_db_path(normalized)
        
        return {
            "original": folder,
            "normalized": normalized,
            "exists": exists,
            "is_dir": is_dir,
            "db_path": db_path,
            "os_info": {
                "platform": sys.platform,
                "cwd": os.getcwd(),
                "path_sep": os.path.sep
            }
        }
    except Exception as e:
        error_msg = str(e)
        trace = traceback.format_exc()
        logger.error(f"调试路径出错: {error_msg}\n{trace}")
        return {"error": error_msg, "trace": trace}

# API路由：启动文件监控
@app.post("/start-monitoring")
async def api_start_monitoring(folder_req: FolderRequest):
    folder = folder_req.folder
    logger.info(f"收到启动文件监控请求: {folder}")
    
    try:
        folder = normalize_path(folder)
        
        # 检查路径是否存在
        if not os.path.exists(folder):
            error = f"路径不存在: {folder}"
            logger.error(error)
            return {"success": False, "message": error}
        
        # 检查是否是文件夹
        if not os.path.isdir(folder):
            error = f"路径不是一个有效的文件夹: {folder}"
            logger.error(error)
            return {"success": False, "message": error}
        
        # 检查是否已创建索引
        db_path = get_db_path(folder)
        if not os.path.exists(db_path) or not os.path.exists(os.path.join(db_path, "index.faiss")):
            error = f"索引不存在，请先创建索引"
            logger.error(error)
            return {"success": False, "message": error}
        
        # 启动监控
        result = start_file_monitoring(folder)
        if result:
            return {"success": True, "message": "文件监控已启动"}
        else:
            return {"success": False, "message": "启动文件监控失败"}
    except Exception as e:
        error_msg = str(e)
        logger.error(f"启动文件监控时出错: {error_msg}")
        return {"success": False, "message": f"启动文件监控时出错: {error_msg}"}

# API路由：停止文件监控
@app.post("/stop-monitoring")
async def api_stop_monitoring(folder_req: FolderRequest):
    folder = folder_req.folder
    logger.info(f"收到停止文件监控请求: {folder}")
    
    try:
        folder = normalize_path(folder)
        result = stop_file_monitoring(folder)
        if result:
            return {"success": True, "message": "文件监控已停止"}
        else:
            return {"success": False, "message": "该文件夹未在监控中"}
    except Exception as e:
        error_msg = str(e)
        logger.error(f"停止文件监控时出错: {error_msg}")
        return {"success": False, "message": f"停止文件监控时出错: {error_msg}"}

# API路由：获取监控状态
@app.get("/monitoring-status")
async def api_monitoring_status(folder: str = None):
    try:
        if folder:
            # 解码URL编码的路径
            folder = urllib.parse.unquote(folder)
            # 规范化路径
            folder = normalize_path(folder)
            # 检查特定文件夹
            is_monitoring = folder in active_observers
            return {"success": True, "is_monitoring": is_monitoring}
        else:
            # 返回所有被监控的文件夹
            monitored_folders = list(active_observers.keys())
            return {"success": True, "monitored_folders": monitored_folders}
    except Exception as e:
        error_msg = str(e)
        logger.error(f"获取监控状态时出错: {error_msg}")
        return {"success": False, "message": f"获取监控状态时出错: {error_msg}"}

# API路由：停止所有监控
@app.post("/stop-all-monitoring")
async def api_stop_all_monitoring():
    """停止所有文件监控"""
    logger.info("请求停止所有文件监控")
    try:
        # 调用停止所有监控函数
        stop_all_monitoring()
        return {"success": True, "message": "已停止所有文件监控"}
    except Exception as e:
        error_msg = str(e)
        logger.error(f"停止所有文件监控时出错: {error_msg}")
        return {"success": False, "message": f"停止所有文件监控时出错: {error_msg}"}

# API路由：清理所有索引数据
@app.post("/clean-all-indexes")
async def clean_all_indexes():
    """清理所有索引数据"""
    logger.info("请求清理所有索引数据")
    try:
        # 停止所有文件监控
        stop_all_monitoring()
        
        # 获取向量存储根目录
        vector_store_dir = os.path.join(os.path.dirname(os.path.dirname(os.path.abspath(__file__))), "vector_store")
        
        # 如果目录存在，删除其中的所有文件和子目录
        if os.path.exists(vector_store_dir):
            # 遍历并删除子目录和文件
            for item in os.listdir(vector_store_dir):
                item_path = os.path.join(vector_store_dir, item)
                if os.path.isdir(item_path):
                    shutil.rmtree(item_path)
                else:
                    os.remove(item_path)
            
            # 保留vector_store目录本身
            logger.info(f"已清理所有索引数据，保留根目录: {vector_store_dir}")
        else:
            # 如果目录不存在，创建它
            os.makedirs(vector_store_dir, exist_ok=True)
            logger.info(f"索引数据目录不存在，已创建: {vector_store_dir}")
        
        return {"success": True, "message": "已清理所有索引数据"}
    except Exception as e:
        error_msg = str(e)
        logger.error(f"清理所有索引数据时出错: {error_msg}")
        return {"success": False, "message": f"清理所有索引数据时出错: {error_msg}"}

# API路由：获取或更新配置
@app.post("/config")
async def update_config(config_req: ConfigRequest):
    """更新系统配置"""
    global MAX_TEXT_LENGTH, MAX_CHUNK_COUNT, MAX_FILE_SIZE_MB, EMBEDDING_MODEL_NAME
    
    try:
        # 检查并更新每个配置项
        changes = []
        
        if config_req.max_text_length is not None:
            if config_req.max_text_length >= 1000 and config_req.max_text_length <= 100000:
                old_value = MAX_TEXT_LENGTH
                MAX_TEXT_LENGTH = config_req.max_text_length
                changes.append(f"最大文本长度: {old_value} -> {MAX_TEXT_LENGTH}")
            else:
                return {"success": False, "message": "最大文本长度必须在1000到100000之间"}
        
        if config_req.max_chunk_count is not None:
            if config_req.max_chunk_count >= 50 and config_req.max_chunk_count <= 500:
                old_value = MAX_CHUNK_COUNT
                MAX_CHUNK_COUNT = config_req.max_chunk_count
                changes.append(f"最大分块数量: {old_value} -> {MAX_CHUNK_COUNT}")
            else:
                return {"success": False, "message": "最大分块数量必须在50到500之间"}
        
        if config_req.max_file_size_mb is not None:
            if config_req.max_file_size_mb >= 10 and config_req.max_file_size_mb <= 500:
                old_value = MAX_FILE_SIZE_MB
                MAX_FILE_SIZE_MB = config_req.max_file_size_mb
                changes.append(f"最大文件大小(MB): {old_value} -> {MAX_FILE_SIZE_MB}")
            else:
                return {"success": False, "message": "最大文件大小必须在10MB到500MB之间"}
        
        if config_req.embedding_model is not None:
            if config_req.embedding_model in EMBEDDING_MODELS:
                old_value = EMBEDDING_MODEL_NAME
                EMBEDDING_MODEL_NAME = config_req.embedding_model
                # 更新token限制
                update_token_limit()
                changes.append(f"嵌入模型: {old_value} -> {EMBEDDING_MODEL_NAME} (Token限制: {MAX_TEXT_BLOCK_SIZE})")
            else:
                return {"success": False, "message": f"不支持的嵌入模型: {config_req.embedding_model}"}
        
        # 记录更改
        if changes:
            logger.info(f"配置已更新: {', '.join(changes)}")
        
        # 返回更新后的配置
        return {
            "success": True, 
            "message": "配置已更新", 
            "config": {
                "max_text_length": MAX_TEXT_LENGTH,
                "max_chunk_count": MAX_CHUNK_COUNT,
                "max_file_size_mb": MAX_FILE_SIZE_MB,
                "embedding_model": {
                    "name": EMBEDDING_MODEL_NAME,
                    "max_tokens": MAX_TEXT_BLOCK_SIZE,
                    "details": EMBEDDING_MODELS.get(EMBEDDING_MODEL_NAME, {})
                },
                "available_models": EMBEDDING_MODELS
            }
        }
    except Exception as e:
        error_msg = str(e)
        logger.error(f"更新配置时出错: {error_msg}")
        return {"success": False, "message": f"更新配置时出错: {error_msg}"}

@app.get("/config")
async def get_config():
    """获取系统配置"""
    return {
        "max_text_length": MAX_TEXT_LENGTH,
        "max_chunk_count": MAX_CHUNK_COUNT,
        "max_file_size_mb": MAX_FILE_SIZE_MB,
        "embedding_model": {
            "name": EMBEDDING_MODEL_NAME,
            "max_tokens": MAX_TEXT_BLOCK_SIZE,
            "details": EMBEDDING_MODELS.get(EMBEDDING_MODEL_NAME, {})
        },
        "available_models": EMBEDDING_MODELS
    }

# 主入口点
if __name__ == "__main__":
    import time
    start_time = time.time()
    print(f"[{time.strftime('%Y-%m-%d %H:%M:%S')}] Python后端启动中...")
    
    # 打印工作目录信息
    print(f"[{time.strftime('%Y-%m-%d %H:%M:%S')}] 当前工作目录: {os.getcwd()}")
    
    # 注册退出处理函数，确保关闭所有监控器
    import atexit
    def cleanup():
        print(f"[{time.strftime('%Y-%m-%d %H:%M:%S')}] 程序退出，清理资源...")
        # 停止所有监控
        stop_all_monitoring()
    atexit.register(cleanup)
    
    # 发送一个特殊标记，表示Python后端已经准备好启动服务器
    # 这个标记会被Electron捕获，用来立即通知前端
    print("PYTHON_BACKEND_READY", flush=True)
    print(f"[{time.strftime('%Y-%m-%d %H:%M:%S')}] Python后端就绪标记已发送")
    
    # 把监控器停止和配置加载放到后台线程中，避免阻塞主启动流程
    def init_background_tasks():
        # 确保没有遗留的监控器
        try:
            thread_start_time = time.time()
            print(f"[{time.strftime('%Y-%m-%d %H:%M:%S')}] 停止所有可能遗留的监控器...")
            stop_all_monitoring()
            print(f"[{time.strftime('%Y-%m-%d %H:%M:%S')}] 停止监控器完成，耗时: {time.time() - thread_start_time:.3f}秒")
            
            # 尝试加载应用配置
            print(f"[{time.strftime('%Y-%m-%d %H:%M:%S')}] 开始加载应用配置...")
            config_load_start = time.time()
            config_dir = os.path.join(os.path.dirname(os.path.dirname(os.path.abspath(__file__))), "vector_store")
            config_file = os.path.join(config_dir, "app_config.json")
            if os.path.exists(config_file):
                print(f"[{time.strftime('%Y-%m-%d %H:%M:%S')}] 发现应用配置文件，加载配置: {config_file}")
                with open(config_file, 'r', encoding='utf-8') as f:
                    config = json.load(f)
                    
                # 加载文本长度限制
                if "MAX_TEXT_LENGTH" in config:
                    global MAX_TEXT_LENGTH
                    MAX_TEXT_LENGTH = config["MAX_TEXT_LENGTH"]
                    print(f"[{time.strftime('%Y-%m-%d %H:%M:%S')}] 已加载文本长度限制: {MAX_TEXT_LENGTH}")
                    
                # 加载块数量限制
                if "MAX_CHUNK_COUNT" in config:
                    global MAX_CHUNK_COUNT
                    MAX_CHUNK_COUNT = config["MAX_CHUNK_COUNT"]
                    print(f"[{time.strftime('%Y-%m-%d %H:%M:%S')}] 已加载块数量限制: {MAX_CHUNK_COUNT}")
            else:
                print(f"[{time.strftime('%Y-%m-%d %H:%M:%S')}] 应用配置文件不存在: {config_file}")
            print(f"[{time.strftime('%Y-%m-%d %H:%M:%S')}] 加载应用配置完成，耗时: {time.time() - config_load_start:.3f}秒")
            
            # 尝试恢复之前的监控状态
            print(f"[{time.strftime('%Y-%m-%d %H:%M:%S')}] 开始恢复监控状态...")
            monitor_start = time.time()
            # 检查是否有保存的配置文件
            config_file = os.path.join(config_dir, "monitoring_config.json")
            if os.path.exists(config_file):
                print(f"[{time.strftime('%Y-%m-%d %H:%M:%S')}] 发现监控配置文件: {config_file}，尝试恢复监控状态")
                with open(config_file, 'r', encoding='utf-8') as f:
                    monitoring_config = json.load(f)
                    
                # 恢复每个目录的监控状态
                for folder, status in monitoring_config.items():
                    if status.get("monitoring", False) and os.path.exists(folder):
                        print(f"[{time.strftime('%Y-%m-%d %H:%M:%S')}] 正在恢复对 {folder} 的监控")
                        inner_start = time.time()
                        start_file_monitoring(folder)
                        print(f"[{time.strftime('%Y-%m-%d %H:%M:%S')}] 恢复对 {folder} 的监控完成，耗时: {time.time() - inner_start:.3f}秒")
            else:
                print(f"[{time.strftime('%Y-%m-%d %H:%M:%S')}] 监控配置文件不存在: {config_file}")
            print(f"[{time.strftime('%Y-%m-%d %H:%M:%S')}] 恢复监控状态完成，耗时: {time.time() - monitor_start:.3f}秒")
            print(f"[{time.strftime('%Y-%m-%d %H:%M:%S')}] 后台初始化任务全部完成，总耗时: {time.time() - thread_start_time:.3f}秒")
        except Exception as e:
            print(f"[{time.strftime('%Y-%m-%d %H:%M:%S')}] 后台初始化任务出错: {str(e)}")
            print(f"[{time.strftime('%Y-%m-%d %H:%M:%S')}] 错误详情: {traceback.format_exc()}")
    
    # 增加日志，记录服务器启动
    print(f"[{time.strftime('%Y-%m-%d %H:%M:%S')}] 准备启动FastAPI服务器...")
    
    # 创建后台线程进行初始化
    import threading
    init_thread = threading.Thread(target=init_background_tasks, daemon=True)
    init_thread.start()
    print(f"[{time.strftime('%Y-%m-%d %H:%M:%S')}] 后台初始化线程已启动")
    
    # 输出总启动时间
    print(f"[{time.strftime('%Y-%m-%d %H:%M:%S')}] 启动预处理完成，耗时: {time.time() - start_time:.3f}秒")
    
    # 启动服务器
    uvicorn.run(app, host="127.0.0.1", port=8000)